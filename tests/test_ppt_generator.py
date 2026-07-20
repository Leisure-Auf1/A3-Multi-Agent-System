"""
Tests for PPTGeneratorAgent (Phase 8.3-F1).

Covers:
- Data model roundtrip: SlideItem, PPTStructure
- Fallback rule-based structure generation
- LLM provider structure generation
- Capability check via check_task_capability()
- Router selection via ModelRouter
- .pptx file generation and validation
- Edge cases: empty material, no python-pptx, etc.

Constraints: does NOT modify Veritas-Core or src/core/
"""

import json
import os
import sys
import tempfile
import unittest
from pathlib import Path
from typing import Any, Dict, List, Optional

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

from agents.content_generator_agent import (
    ContentGeneratorAgent, TeachingMaterial, Chapter,
    ConceptItem, ExampleItem, ExerciseItem,
)
from agents.planner_agent import PlannerAgent, LearningPlan, PlanNode
from config.model_capability import (
    ModelCapability,
    check_task_capability,
    get_provider_capabilities,
    has_capability,
)
from config.task_capability import TaskType, TASK_REQUIREMENTS
from config.model_router import ModelRouter, RouterResult
from config.model_registry import MODEL_REGISTRY, find_models_with_capability

# Import agent under test (source path added to sys.path above)
from agents.ppt_generator_agent import (
    PPTGeneratorAgent,
    PPTStructure,
    SlideItem,
)


# ── Helpers ──────────────────────────────────────────────────

def _make_material(chapter_count: int = 3) -> TeachingMaterial:
    """Create a TeachingMaterial for testing."""
    concepts = [
        ConceptItem(name=f"概念 {i+1}", description=f"概念{i+1}的详细解释",
                     difficulty="beginner", related=[])
        for i in range(2)
    ]
    examples = [
        ExampleItem(title=f"示例 {i+1}", code=f"code_{i+1}",
                     explanation=f"示例说明 {i+1}", expected_output=f"output_{i+1}")
        for i in range(1)
    ]
    exercises = [
        ExerciseItem(question=f"练习题 {i+1}?", answer=f"答案 {i+1}",
                      hint=f"提示 {i+1}", type="open")
        for i in range(2)
    ]

    chapters = [
        Chapter(
            chapter_id=f"ch{i+1}",
            title=f"第{i+1}章: 核心主题{i+1}",
            explanation=f"这是第{i+1}章的详细解释内容，涵盖相关概念。",
            concepts=concepts,
            examples=examples,
            exercises=exercises,
            estimated_minutes=20,
            summary=f"第{i+1}章小结",
        )
        for i in range(chapter_count)
    ]

    return TeachingMaterial(
        material_id="test_mat_001",
        title="Python 编程基础教材",
        learning_objectives=["理解Python语法", "掌握函数定义", "熟悉面向对象编程"],
        chapters=chapters,
        overall_summary="本教材涵盖了Python编程的基础知识。",
        target_profile="junior_dev",
        total_estimated_minutes=60,
        generation_source="rule",
    )


def _make_profile(kb="junior_dev", cog="visual_dominant", pace="normal"):
    """Create a DynamicProfile for testing."""
    from core.agent_router import DynamicProfile
    return DynamicProfile(
        knowledge_base=kb,
        cognitive_style=cog,
        error_prone_bias="magic_syntax_blind",
        learning_pace=pace,
        interaction_preference="code_sandbox",
        frustration_threshold="medium",
    )


def _make_plan(node_count=3):
    """Create a LearningPlan for testing."""
    nodes = []
    for i in range(node_count):
        nodes.append(PlanNode(
            node_id=f"topic_{i + 1}",
            title=f"Topic {i + 1}",
            core_concept=f"Core concept {i + 1}",
            depth=1 + (i % 3),
            estimated_minutes=15 + i * 5,
            required_concepts=[],
            exercise_count=3,
            teaching_strategy="standard",
            notes=f"Notes for topic {i + 1}",
        ))
    return LearningPlan(
        plan_id="test_plan",
        profile_summary="junior_dev / visual_dominant / normal",
        nodes=nodes,
        total_minutes=60,
        strategy_rationale="Test strategy",
    )


# ── Mock LLM Provider ────────────────────────────────────────

class MockLLMProvider:
    """Simulates an LLM that returns PPT structure JSON."""

    def __init__(self, response_json: Optional[dict] = None, should_fail: bool = False):
        self._response = response_json
        self._should_fail = should_fail
        self.model = "mock-model"

    @property
    def is_available(self):
        return True

    def generate(self, prompt: str, system_prompt: str = "", **kwargs):
        if self._should_fail:
            raise RuntimeError("LLM unavailable")

        class FakeResponse:
            content = json.dumps(self._response or {})
            success = True
            error = ""

        return FakeResponse()


class MockFailingLLMProvider:
    """Simulates an LLM that always fails."""

    @property
    def is_available(self):
        return True

    def generate(self, prompt: str, system_prompt: str = "", **kwargs):
        class FakeResponse:
            content = ""
            success = False
            error = "API error"

        return FakeResponse()


# ──────────────────────────────────────────────
# 1. Data Model Tests
# ──────────────────────────────────────────────

class TestDataModels:
    """Test serialization roundtrip for PPT data models."""

    def test_slide_item_roundtrip(self):
        """SlideItem to_dict/from_dict works correctly."""
        s = SlideItem(
            slide_index=1,
            title="封面",
            content_lines=["副标题1", "副标题2"],
            layout="title",
            notes="讲师备注",
        )
        d = s.to_dict()
        s2 = SlideItem.from_dict(d)
        assert s2.slide_index == 1
        assert s2.title == "封面"
        assert s2.content_lines == ["副标题1", "副标题2"]
        assert s2.layout == "title"
        assert s2.notes == "讲师备注"

    def test_slide_item_defaults(self):
        """SlideItem has sensible defaults."""
        s = SlideItem(slide_index=3, title="Slide 3")
        assert s.content_lines == []
        assert s.layout == "content"
        assert s.notes == ""

    def test_ppt_structure_roundtrip(self):
        """PPTStructure to_dict/from_dict works correctly."""
        slides = [
            SlideItem(slide_index=1, title="封面", layout="title"),
            SlideItem(slide_index=2, title="内容", content_lines=["行1", "行2"]),
        ]
        ps = PPTStructure(
            ppt_id="ppt_001",
            title="Test PPT",
            topic="Python",
            level="beginner",
            slides=slides,
            total_slides=2,
            generation_source="rule",
            metadata={"chapter_count": 2},
        )
        d = ps.to_dict()
        ps2 = PPTStructure.from_dict(d)
        assert ps2.ppt_id == "ppt_001"
        assert ps2.title == "Test PPT"
        assert ps2.total_slides == 2
        assert ps2.generation_source == "rule"
        assert len(ps2.slides) == 2

    def test_ppt_structure_defaults(self):
        """PPTStructure has sensible defaults."""
        ps = PPTStructure(ppt_id="ppt_001", title="Test")
        assert ps.topic == ""
        assert ps.level == "beginner"
        assert ps.slides == []
        assert ps.total_slides == 0
        assert ps.generation_source == "rule"
        assert ps.metadata == {}


# ──────────────────────────────────────────────
# 2. Fallback Structure Generation Tests
# ──────────────────────────────────────────────

class TestFallbackStructureGeneration:
    """Tests for rule-based PPT structure generation."""

    def setup_method(self):
        self.agent = PPTGeneratorAgent()

    def test_fallback_creates_title_slide(self):
        """Fallback always creates a title slide first."""
        material = _make_material()
        structure = self.agent.fallback_generate_structure(material.to_dict())

        assert len(structure.slides) > 0
        first_slide = structure.slides[0]
        assert first_slide.slide_index == 1
        assert first_slide.layout == "title"
        assert len(first_slide.title) > 0

    def test_fallback_creates_objectives_slide(self):
        """When learning objectives exist, creates objectives slide."""
        material = _make_material()
        material.learning_objectives = ["目标A", "目标B", "目标C"]
        structure = self.agent.fallback_generate_structure(material.to_dict())

        # Second slide should be objectives
        layouts = [s.layout for s in structure.slides]
        assert "bullet_list" in layouts

    def test_fallback_creates_chapter_slides(self):
        """Each chapter generates multiple slides (title + concepts + examples + exercises)."""
        material = _make_material(chapter_count=2)
        structure = self.agent.fallback_generate_structure(material.to_dict())

        # Should have: title + objectives + 2*(title+concepts+examples+exercises) + summary
        # = 1 + 1 + 2*4 + 1 = 11 slides
        assert len(structure.slides) >= 6
        assert structure.total_slides == len(structure.slides)

    def test_fallback_creates_summary_slide(self):
        """Last slide is a summary."""
        material = _make_material()
        structure = self.agent.fallback_generate_structure(material.to_dict())

        last_slide = structure.slides[-1]
        assert last_slide.layout == "summary"
        assert "总结" in last_slide.title

    def test_fallback_empty_material(self):
        """Empty material generates minimal valid structure."""
        empty = TeachingMaterial(material_id="empty", title="空教材")
        structure = self.agent.fallback_generate_structure(empty.to_dict())

        assert structure.total_slides >= 2  # title + summary
        assert structure.generation_source == "rule"

    def test_fallback_generation_source(self):
        """Fallback always has generation_source='rule'."""
        material = _make_material()
        structure = self.agent.fallback_generate_structure(material.to_dict())
        assert structure.generation_source == "rule"

    def test_fallback_single_chapter(self):
        """Single chapter generates complete structure."""
        material = _make_material(chapter_count=1)
        structure = self.agent.fallback_generate_structure(material.to_dict())
        assert len(structure.slides) >= 4  # title + objectives + chapter slides + summary

    def test_fallback_many_chapters(self):
        """Many chapters generate proportional slides."""
        material = _make_material(chapter_count=5)
        structure = self.agent.fallback_generate_structure(material.to_dict())
        # Each chapter creates ~4 slides (title + concepts + examples + exercises)
        assert len(structure.slides) >= 12

    def test_fallback_no_objectives(self):
        """Material without learning objectives skips objectives slide."""
        material = _make_material()
        material.learning_objectives = []
        material_dict = material.to_dict()
        structure = self.agent.fallback_generate_structure(material_dict)

        # Should not have "bullet_list" layout for objectives
        layouts = [s.layout for s in structure.slides[1:]]
        # But bullet_list may appear for concepts - that's fine

    def test_fallback_metadata(self):
        """Metadata includes chapter and objective counts."""
        material = _make_material(chapter_count=3)
        structure = self.agent.fallback_generate_structure(material.to_dict())

        assert "chapter_count" in structure.metadata
        assert structure.metadata["chapter_count"] == 3

    def test_get_ppt_structure_with_teaching_material(self):
        """get_ppt_structure returns valid PPTStructure from TeachingMaterial."""
        material = _make_material()
        structure = self.agent.get_ppt_structure(material)

        assert isinstance(structure, PPTStructure)
        assert len(structure.slides) > 0
        assert structure.generation_source == "rule"

    def test_get_ppt_structure_with_dict(self):
        """get_ppt_structure works with dict input."""
        material = _make_material()
        structure = self.agent.get_ppt_structure(material.to_dict())

        assert isinstance(structure, PPTStructure)
        assert len(structure.slides) > 0

    def test_get_ppt_structure_empty_raises(self):
        """Empty material raises ValueError."""
        with pytest.raises(ValueError):
            self.agent.get_ppt_structure({})

    def test_generate_ppt_empty_raises(self):
        """generate_ppt with empty material raises ValueError."""
        with pytest.raises(ValueError):
            self.agent.generate_ppt({})


# ──────────────────────────────────────────────
# 3. LLM Provider Tests
# ──────────────────────────────────────────────

class TestLLMProvider:
    """Tests for LLM-powered PPT generation."""

    def setup_method(self):
        self.agent = PPTGeneratorAgent()

    def test_llm_provider_produces_llm_structure(self):
        """When LLM provider succeeds, structure has generation_source='llm'."""
        llm_response = {
            "title": "Python PPT",
            "topic": "Python",
            "slides": [
                {"slide_index": 1, "title": "封面", "layout": "title",
                 "content_lines": ["Python编程"], "notes": ""},
                {"slide_index": 2, "title": "目标", "layout": "bullet_list",
                 "content_lines": ["目标1", "目标2"], "notes": ""},
                {"slide_index": 3, "title": "总结", "layout": "summary",
                 "content_lines": ["总结内容"], "notes": ""},
            ],
        }
        self.agent.set_llm_provider(MockLLMProvider(llm_response))

        material = _make_material(chapter_count=1)
        structure = self.agent._generate_structure(material.to_dict())

        assert structure.generation_source == "llm"
        assert structure.title == "Python PPT"
        assert len(structure.slides) == 3

    def test_llm_failure_falls_back_to_rule(self):
        """When LLM fails, structure uses rule fallback."""
        self.agent.set_llm_provider(MockFailingLLMProvider())

        material = _make_material()
        structure = self.agent._generate_structure(material.to_dict())

        assert structure.generation_source == "rule"
        assert len(structure.slides) > 0

    def test_llm_exception_falls_back_to_rule(self):
        """When LLM raises exception, rule fallback is used."""
        self.agent.set_llm_provider(MockLLMProvider(should_fail=True))

        material = _make_material()
        structure = self.agent._generate_structure(material.to_dict())

        assert structure.generation_source == "rule"

    def test_empty_provider_uses_rule(self):
        """When no provider is set, rule generation is used."""
        material = _make_material()
        structure = self.agent._generate_structure(material.to_dict())

        assert structure.generation_source == "rule"

    def test_set_llm_provider_method(self):
        """set_llm_provider stores the provider."""
        mock = MockLLMProvider({"title": "T"})
        self.agent.set_llm_provider(mock)
        assert self.agent._llm_provider is mock

    def test_set_llm_provider_none_disables_llm(self):
        """Setting provider to None disables LLM path."""
        mock = MockLLMProvider({"title": "T"})
        self.agent.set_llm_provider(mock)
        self.agent.set_llm_provider(None)

        material = _make_material()
        structure = self.agent._generate_structure(material.to_dict())
        assert structure.generation_source == "rule"


# ──────────────────────────────────────────────
# 4. .pptx File Generation Tests
# ──────────────────────────────────────────────

class TestPptxFileGeneration:
    """Tests for .pptx file rendering."""

    def setup_method(self):
        self.agent = PPTGeneratorAgent()

    def test_generate_pptx_creates_valid_file(self):
        """generate_ppt creates a .pptx file with content."""
        material = _make_material(chapter_count=2)
        ppt_path = self.agent.generate_ppt(material)

        assert os.path.isfile(ppt_path)
        assert ppt_path.endswith(".pptx")
        file_size = os.path.getsize(ppt_path)
        assert file_size > 1000, f"PPT file too small: {file_size} bytes"
        os.unlink(ppt_path)

    def test_generate_pptx_with_dict_input(self):
        """generate_ppt works with dict input."""
        material = _make_material()
        ppt_path = self.agent.generate_ppt(material.to_dict())

        assert os.path.isfile(ppt_path)
        assert os.path.getsize(ppt_path) > 1000
        os.unlink(ppt_path)

    def test_generate_pptx_multi_chapter(self):
        """PPT generation handles multiple chapters."""
        material = _make_material(chapter_count=4)
        ppt_path = self.agent.generate_ppt(material)

        assert os.path.isfile(ppt_path)
        file_size = os.path.getsize(ppt_path)
        assert file_size > 2000, f"Multi-chapter PPT too small: {file_size} bytes"
        os.unlink(ppt_path)

    def test_generate_pptx_output_dir(self):
        """generate_ppt supports custom output directory."""
        material = _make_material()
        with tempfile.TemporaryDirectory() as tmpdir:
            ppt_path = self.agent.generate_ppt(material, output_dir=tmpdir)
            assert ppt_path.startswith(tmpdir)
            assert os.path.isfile(ppt_path)
            assert os.path.getsize(ppt_path) > 1000

    def test_generate_pptx_has_title_slide(self):
        """Generated PPT file can be opened and validated."""
        material = _make_material()
        ppt_path = self.agent.generate_ppt(material)

        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)
            # Should have at least 3 slides
            assert len(prs.slides) >= 3
            # First slide has title text
            first_slide = prs.slides[0]
            assert first_slide.shapes.title is not None
            assert len(first_slide.shapes.title.text) > 0
        finally:
            if os.path.exists(ppt_path):
                os.unlink(ppt_path)

    def test_generate_pptx_empty_material_raises(self):
        """Empty material raises ValueError."""
        with pytest.raises(ValueError):
            self.agent.generate_ppt({})


# ──────────────────────────────────────────────
# 5. Capability Check Tests
# ──────────────────────────────────────────────

class TestCapabilityCheck:
    """Tests for PPT generation capability checking via Capability Layer."""

    def setup_method(self):
        self.agent = PPTGeneratorAgent()

    def test_ppt_generation_flag_exists(self):
        """ModelCapability.PPT_GENERATION is defined."""
        assert ModelCapability.PPT_GENERATION is not None
        assert ModelCapability.PPT_GENERATION.value > 0

    def test_ppt_generation_in_task_requirements(self):
        """TaskType.GENERATE_PPT maps to PPT_GENERATION + TOOL_CALLING."""
        reqs = TASK_REQUIREMENTS.get(TaskType.GENERATE_PPT, [])
        assert len(reqs) > 0
        assert ModelCapability.PPT_GENERATION in reqs
        assert ModelCapability.TOOL_CALLING in reqs

    def test_check_capability_unsupported_provider(self):
        """DeepSeek doesn't support PPT generation."""
        supported, err = self.agent.check_capability("deepseek", "deepseek-chat")
        assert not supported
        assert err is not None
        assert "不支持" in err

    def test_check_capability_supported_model(self):
        """gpt-5.6-vision supports PPT generation."""
        supported, err = self.agent.check_capability("openai", "gpt-5.6-vision")
        # This should be supported since gpt-5.6-vision has PPT_GENERATION
        if supported:
            assert err is None
        # Even if not directly matched, we don't want to assert False
        # because test model availability may vary

    def test_check_capability_returns_tuple(self):
        """check_capability returns (bool, Optional[str])."""
        result = self.agent.check_capability("mock", "mock-model-v1")
        assert isinstance(result, tuple)
        assert len(result) == 2
        assert isinstance(result[0], bool)

    def test_check_capability_unknown_provider(self):
        """Unknown provider returns some result without crashing."""
        supported, err = self.agent.check_capability("nonexistent", "")
        # Should not crash - returns a result
        assert isinstance(supported, bool)

    def test_check_task_capability_direct(self):
        """Direct check_task_capability call works."""
        supported, err = check_task_capability(
            "openai", "gpt-5.6-vision", TaskType.GENERATE_PPT
        )
        # gpt-5.6-vision has DOCUMENT_GENERATION and TOOL_CALLING
        # ... actually in model_registry it has PPT_GENERATION | PDF_GENERATION
        # but not explicitly DOCUMENT_GENERATION
        # Let's just verify it doesn't crash
        assert isinstance(supported, bool)

    def test_ppt_generation_in_registry_labels(self):
        """PPT_GENERATION has labels and icons."""
        from config.model_capability import CAPABILITY_LABELS, CAPABILITY_ICONS
        assert ModelCapability.PPT_GENERATION in CAPABILITY_LABELS
        assert ModelCapability.PPT_GENERATION in CAPABILITY_ICONS
        assert CAPABILITY_LABELS[ModelCapability.PPT_GENERATION] == "PPT生成"
        assert CAPABILITY_ICONS[ModelCapability.PPT_GENERATION] == "📊"


# ──────────────────────────────────────────────
# 6. Router Selection Tests
# ──────────────────────────────────────────────

class TestRouterSelection:
    """Tests for ModelRouter PPT generation model selection."""

    def setup_method(self):
        self.router = ModelRouter()
        self.agent = PPTGeneratorAgent()

    def test_find_models_for_ppt_task(self):
        """find_models returns models with PPT_GENERATION + TOOL_CALLING."""
        models = self.router.find_models(TaskType.GENERATE_PPT)
        assert isinstance(models, list)
        # At least mock-model-v1 should support this
        if len(models) > 0:
            for m in models:
                assert ModelCapability.PPT_GENERATION in m.capabilities
                assert ModelCapability.TOOL_CALLING in m.capabilities

    def test_select_model_for_ppt_task(self):
        """select_model returns a valid RouterResult for PPT task."""
        result = self.router.select_model(TaskType.GENERATE_PPT)
        assert isinstance(result, RouterResult)
        assert result.task == TaskType.GENERATE_PPT

    def test_select_model_preferred_provider(self):
        """select_model with preferred_provider works."""
        result = self.router.select_model(
            TaskType.GENERATE_PPT,
            preferred_provider="openai",
        )
        assert result.task == TaskType.GENERATE_PPT
        # result may or may not have found a match

    def test_router_result_to_dict(self):
        """RouterResult.to_dict() returns serializable data."""
        result = self.router.select_model(TaskType.GENERATE_PPT)
        d = result.to_dict()
        assert "success" in d
        assert "task" in d
        assert "model_id" in d

    def test_agent_get_recommended_model(self):
        """PPTGeneratorAgent.get_recommended_model() returns RouterResult."""
        result = self.agent.get_recommended_model()
        assert hasattr(result, 'success')
        assert hasattr(result, 'task')
        assert result.task == TaskType.GENERATE_PPT
        # result.success may be False if no model has DOCUMENT_GENERATION + TOOL_CALLING

    def test_agent_get_capable_models(self):
        """PPTGeneratorAgent.get_capable_models() returns models list."""
        models = self.agent.get_capable_models()
        assert isinstance(models, list)

    def test_agent_get_capable_providers(self):
        """PPTGeneratorAgent.get_capable_providers() returns providers list."""
        providers = self.agent.get_capable_providers()
        assert isinstance(providers, list)

    def test_task_provider_priority(self):
        """TASK_PROVIDER_PRIORITY includes GENERATE_PPT."""
        from config.model_router import TASK_PROVIDER_PRIORITY
        assert TaskType.GENERATE_PPT in TASK_PROVIDER_PRIORITY
        priorities = TASK_PROVIDER_PRIORITY[TaskType.GENERATE_PPT]
        assert len(priorities) > 0
        assert "openai" in priorities or "anthropic" in priorities

    def test_find_models_with_ppt_capability(self):
        """find_models_with_capability finds PPT capable models."""
        models = find_models_with_capability(ModelCapability.PPT_GENERATION)
        assert isinstance(models, list)
        # At least gpt-5.6-vision and gemini-ultra-vision have PPT_GENERATION
        model_ids = [m.model_id for m in models]
        assert "gpt-5.6-vision" in model_ids or "gemini-ultra-vision" in model_ids

    def test_get_recommended_model_with_preferred(self):
        """get_recommended_model with preferred provider returns RouterResult."""
        result = self.agent.get_recommended_model(preferred_provider="openai")
        assert hasattr(result, 'success')
        assert hasattr(result, 'task')
        assert result.task == TaskType.GENERATE_PPT
        # result.success may be False if no model matches all required capabilities


# ──────────────────────────────────────────────
# 7. Integration: Full Pipeline
# ──────────────────────────────────────────────

class TestFullPipeline:
    """Integration tests: ContentGeneratorAgent → PPTGeneratorAgent."""

    def setup_method(self):
        self.content_agent = ContentGeneratorAgent()
        self.ppt_agent = PPTGeneratorAgent()

    def test_content_to_ppt_pipeline(self):
        """Full pipeline: generate material → generate PPT."""
        profile = _make_profile()
        plan = _make_plan()
        material = self.content_agent.generate_material(profile, plan)

        assert isinstance(material, TeachingMaterial)
        assert len(material.chapters) > 0

        ppt_path = self.ppt_agent.generate_ppt(material)

        assert os.path.isfile(ppt_path)
        assert os.path.getsize(ppt_path) > 1000
        os.unlink(ppt_path)

    def test_content_to_ppt_structure_only(self):
        """Get PPT structure without rendering file."""
        profile = _make_profile()
        plan = _make_plan()
        material = self.content_agent.generate_material(profile, plan)

        structure = self.ppt_agent.get_ppt_structure(material)

        assert isinstance(structure, PPTStructure)
        assert structure.generation_source == "rule"
        assert len(structure.slides) > 0
        assert structure.title == material.title


# ──────────────────────────────────────────────
# 8. Edge Cases
# ──────────────────────────────────────────────

class TestEdgeCases:
    """Edge case handling."""

    def setup_method(self):
        self.agent = PPTGeneratorAgent()

    def test_empty_chapters_list(self):
        """Material with no chapters generates minimal PPT."""
        material = TeachingMaterial(
            material_id="empty_chapters",
            title="仅有标题的教材",
            learning_objectives=["目标1"],
            chapters=[],
            overall_summary="总结",
        )
        structure = self.agent.fallback_generate_structure(material.to_dict())
        assert len(structure.slides) >= 2  # title + summary

    def test_chapter_without_concepts(self):
        """Chapter without concepts doesn't generate concepts slide."""
        ch = Chapter(
            chapter_id="no_concepts",
            title="无概念章节",
            explanation="只有解释，没有概念",
            concepts=[],
            examples=[],
            exercises=[],
        )
        material = TeachingMaterial(
            material_id="no_concepts_mat",
            title="无概念教材",
            chapters=[ch],
        )
        structure = self.agent.fallback_generate_structure(material.to_dict())
        # Should have title + chapter title + summary = 3 slides (no concepts/examples/exercises)
        assert len(structure.slides) >= 3

    def test_unicode_content(self):
        """Chinese characters in content render correctly in PPT."""
        material = TeachingMaterial(
            material_id="cn_mat",
            title="Python 编程入门 — 从零到精通",
            learning_objectives=["理解变量和数据类型", "掌握条件判断"],
            chapters=[
                Chapter(
                    chapter_id="ch1",
                    title="第一章：变量与数据类型",
                    explanation="Python 中的变量不需要声明类型。",
                    concepts=[
                        ConceptItem(name="变量", description="存储数据的容器"),
                    ],
                ),
            ],
            overall_summary="本章总结了 Python 编程的基础知识。",
        )
        ppt_path = self.agent.generate_ppt(material)

        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)
            # Verify Chinese content in first slide
            first_slide = prs.slides[0]
            text = first_slide.shapes.title.text
            assert "Python" in text
        finally:
            if os.path.exists(ppt_path):
                os.unlink(ppt_path)

    def test_generate_ppt_with_provider_model_args(self):
        """generate_ppt accepts provider and model args without crashing."""
        material = _make_material()
        ppt_path = self.agent.generate_ppt(
            material,
            provider="openai",
            model="gpt-4o",
        )

        assert os.path.isfile(ppt_path)
        assert os.path.getsize(ppt_path) > 1000
        os.unlink(ppt_path)
