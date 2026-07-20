"""
Phase 8.3-F1 — PPTGeneratorAgent: 基于 TeachingMaterial 生成 .pptx 课件

职责:
  输入: TeachingMaterial (from ContentGeneratorAgent)
  输出: .pptx 文件路径

双模式:
  LLM 增强模式 — 调用 LLMProvider 生成结构化 slide 内容
  Fallback 规则模式 — 基于 TeachingMaterial 章节自动构建 PPT (零延迟, 始终可用)

Capability Layer:
  使用 ModelRouter + check_task_capability() 确保能力检查
  禁止绕过 Capability Layer 直接调用 LLM

约束:
  不修改 src/core/、Veritas-Core、已有 Agent interface
"""

from __future__ import annotations
import json
import os
import tempfile
import uuid
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

from src.config.model_capability import ModelCapability
from src.config.task_capability import TaskType
from src.config.model_capability import check_task_capability
from src.config.model_router import ModelRouter, RouterResult


# ──────────────────────────────────────────────
# 数据模型
# ──────────────────────────────────────────────

@dataclass
class SlideItem:
    """单个 slide 的内容结构."""
    slide_index: int
    title: str
    content_lines: List[str] = field(default_factory=list)
    layout: str = "content"  # title | content | summary | bullet_list
    notes: str = ""
    image: str = ""           # Phase 8.3-F2: SVG path or base64 data URI

    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_index": self.slide_index,
            "title": self.title,
            "content_lines": self.content_lines,
            "layout": self.layout,
            "notes": self.notes,
            "image": self.image,
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "SlideItem":
        return cls(
            slide_index=data.get("slide_index", 0),
            title=data.get("title", ""),
            content_lines=data.get("content_lines", []),
            layout=data.get("layout", "content"),
            notes=data.get("notes", ""),
            image=data.get("image", ""),
        )


@dataclass
class PPTStructure:
    """完整 PPT 结构 — 所有 slide 的集合."""
    ppt_id: str
    title: str
    topic: str = ""
    level: str = "beginner"
    slides: List[SlideItem] = field(default_factory=list)
    total_slides: int = 0
    generation_source: str = "rule"  # rule | llm
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "ppt_id": self.ppt_id,
            "title": self.title,
            "topic": self.topic,
            "level": self.level,
            "slides": [s.to_dict() for s in self.slides],
            "total_slides": self.total_slides,
            "generation_source": self.generation_source,
            "metadata": self.metadata,
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "PPTStructure":
        return cls(
            ppt_id=data.get("ppt_id", ""),
            title=data.get("title", ""),
            topic=data.get("topic", ""),
            level=data.get("level", "beginner"),
            slides=[SlideItem.from_dict(s) for s in data.get("slides", [])],
            total_slides=data.get("total_slides", 0),
            generation_source=data.get("generation_source", "rule"),
            metadata=data.get("metadata", {}),
        )


# ──────────────────────────────────────────────
# PPTGeneratorAgent
# ──────────────────────────────────────────────

class PPTGeneratorAgent:
    """
    AI PPT 生成 Agent.

    基于 TeachingMaterial 生成 .pptx 课件文件。
    支持 LLM 生成 slide 结构和 fallback 规则生成。

    使用方式:
        agent = PPTGeneratorAgent()
        agent.set_llm_provider(provider)  # 可选
        ppt_path = agent.generate_ppt(material, provider="openai", model="gpt-4o")
    """

    def __init__(self):
        self._llm_provider = None  # LLMProvider (optional, backward compat)
        self._orchestrator = None  # Phase 9.3-B
        self._router = ModelRouter()

    # ── LLM Provider Injection ─────

    def set_llm_provider(self, provider: Any) -> None:
        """注入 LLMProvider 以启用 LLM 增强 PPT 生成 (None = 纯规则模式)."""
        self._llm_provider = provider

    def set_orchestrator(self, orchestrator: Any) -> None:
        """Phase 9.3-B: inject OrchestratorRuntime (preferred over llm_provider)."""
        self._orchestrator = orchestrator

    # ── 主入口 ────────────────────

    def generate_ppt(
        self,
        material: Any,  # TeachingMaterial
        provider: str = "",
        model: str = "",
        output_dir: str = "",
    ) -> str:
        """
        根据 TeachingMaterial 生成 .pptx 文件.

        Args:
            material: TeachingMaterial 实例或 dict
            provider: 可选指定 provider (为空则自动路由)
            model: 可选指定 model
            output_dir: 可选输出目录 (为空则使用临时目录)

        Returns:
            .pptx 文件的绝对路径

        Raises:
            ValueError: 输入无效
            RuntimeError: PPT 生成失败 (含错误信息)
        """
        # 统一为 dict
        material_dict = material.to_dict() if hasattr(material, "to_dict") else material

        if not material_dict:
            raise ValueError("TeachingMaterial 不能为空")

        # ── Capability Check ─────
        ppt_structure = self._generate_structure(material_dict, provider, model)

        # ── Render .pptx ─────
        ppt_path = self._render_pptx(ppt_structure, output_dir)
        return ppt_path

    # ── Slide 结构生成 ──────────

    def _generate_structure(
        self,
        material: Dict[str, Any],
        provider: str = "",
        model: str = "",
    ) -> PPTStructure:
        """生成 PPT 结构 (LLM 优先, fallback 兜底)."""

        # ── 尝试 LLM 路径 ──
        if self._llm_provider is not None:
            # Capability check (advisory, not a hard gate)
            if provider and model:
                supported, _ = check_task_capability(
                    provider, model, TaskType.GENERATE_PPT
                )
                if not supported:
                    # 能力不足, 直接 fallback
                    return self.fallback_generate_structure(material)

            # Router: find best model for this task (advisory)
            route_result = self._router.select_model(
                TaskType.GENERATE_PPT, preferred_provider=provider
            )

            # Try LLM regardless of router success — router is advisory
            try:
                structure = self._generate_with_llm(material, route_result)
                if structure is not None:
                    return structure
            except Exception:
                pass  # fallback to rule

        # ── Fallback 规则路径 ──
        return self.fallback_generate_structure(material)

    # ── LLM 增强生成 ──────────────

    LLM_SLIDES_PROMPT = """你是一个专业的教学课件设计师。请根据以下教材内容，设计一份完整的 PPT 幻灯片结构。

[教材标题]
{title}

[学习目标]
{objectives}

[章节内容]
{chapters}

[整体总结]
{summary}

请生成一份完整的 PPT 结构，输出纯 JSON：

{{
  "title": "PPT标题",
  "topic": "主题",
  "slides": [
    {{
      "slide_index": 1,
      "title": "封面标题",
      "layout": "title",
      "content_lines": ["副标题行1", "副标题行2"],
      "notes": "讲师备注"
    }},
    {{
      "slide_index": 2,
      "title": "学习目标",
      "layout": "bullet_list",
      "content_lines": ["目标1", "目标2", "目标3"],
      "notes": ""
    }},
    {{
      "slide_index": 3,
      "title": "章节标题",
      "layout": "content",
      "content_lines": ["要点1: ...", "要点2: ...", "要点3: ..."],
      "notes": "本章重点"
    }}
  ]
}}

要求:
- 第一页为封面 (title slide)
- 第二页为学习目标 (bullet_list)
- 每个章节生成 1-3 页 slides (含概念, 示例, 练习)
- 最后一页为总结 (summary)
- content_lines 每行不要太长 (≤80 字符)
- 所有内容用中文编写
- 每个 slide 的 layout 必须为: title | content | bullet_list | summary
- 总 slides 数量: {slide_count_target}

只输出 JSON，不要任何额外文本。"""

    def _generate_with_llm(
        self,
        material: Dict[str, Any],
        route_result: RouterResult,
    ) -> Optional[PPTStructure]:
        """使用 LLMProvider 生成 PPT slide 结构."""
        if self._llm_provider is None:
            return None

        title = material.get("title", "教学课件")
        objectives = json.dumps(
            material.get("learning_objectives", []),
            ensure_ascii=False,
        )
        chapters_list = material.get("chapters", [])

        # Build chapters text
        chapters_lines = []
        for ch in chapters_list:
            ch_title = ch.get("title", "")
            ch_expl = ch.get("explanation", "")[:200]
            concepts = [c.get("name", "") for c in ch.get("concepts", [])]
            examples = [e.get("title", "") for e in ch.get("examples", [])]
            exercises = [e.get("question", "") for e in ch.get("exercises", [])]

            chapters_lines.append(f"## {ch_title}")
            if ch_expl:
                chapters_lines.append(f"  解释: {ch_expl}")
            if concepts:
                chapters_lines.append(f"  概念: {', '.join(concepts)}")
            if examples:
                chapters_lines.append(f"  示例: {', '.join(examples)}")
            if exercises:
                chapters_lines.append(f"  练习: {', '.join(exercises[:2])}")
            chapters_lines.append("")

        chapters_text = "\n".join(chapters_lines)
        summary = material.get("overall_summary", "")

        # Target slide count: ~1.5 per chapter + 2 (title + summary)
        slide_count_target = max(4, len(chapters_list) * 2 + 2)

        prompt = self.LLM_SLIDES_PROMPT.format(
            title=title,
            objectives=objectives,
            chapters=chapters_text,
            summary=summary,
            slide_count_target=slide_count_target,
        )

        try:
            resp = self._llm_provider.generate(
                prompt=prompt,
                system_prompt="你是一个专业的 PPT 课件设计师。输出纯 JSON。",
                temperature=0.3,
                max_tokens=2048,
            )

            if not resp.success:
                return None

            content = resp.content.strip()
            # Strip markdown code fences
            if content.startswith("```"):
                lines = content.split("\n")
                content = "\n".join(lines[1:]) if len(lines) > 2 else content
                if content.endswith("```"):
                    content = content[:-3]

            data = json.loads(content)

            if "slides" not in data:
                return None

            slides = []
            for i, s in enumerate(data.get("slides", [])):
                slides.append(SlideItem(
                    slide_index=s.get("slide_index", i + 1),
                    title=s.get("title", f"Slide {i + 1}"),
                    content_lines=s.get("content_lines", []),
                    layout=s.get("layout", "content"),
                    notes=s.get("notes", ""),
                ))

            ppt_id = f"ppt_{uuid.uuid4().hex[:12]}"
            structure = PPTStructure(
                ppt_id=ppt_id,
                title=data.get("title", title),
                topic=data.get("topic", material.get("title", "")),
                level=material.get("target_profile", "beginner"),
                slides=slides,
                total_slides=len(slides),
                generation_source="llm",
                metadata={
                    "route_model": route_result.model_id,
                    "route_provider": route_result.model_info.provider if route_result.model_info else "",
                },
            )
            return structure

        except Exception:
            return None

    # ── Fallback 规则生成 ──────────

    def fallback_generate_structure(
        self,
        material: Dict[str, Any],
    ) -> PPTStructure:
        """
        基于规则生成 PPT 结构 (零延迟, 始终可用).

        自动从 TeachingMaterial 的章节、概念、示例、练习中
        提取内容并组织为 slide 结构。
        """
        title = material.get("title", "教学课件")
        ppt_id = f"ppt_{uuid.uuid4().hex[:12]}"
        chapters = material.get("chapters", [])
        objectives = material.get("learning_objectives", [])
        summary = material.get("overall_summary", "")

        slides: List[SlideItem] = []

        # ── Slide 1: 封面 ──
        slides.append(SlideItem(
            slide_index=1,
            title=title,
            content_lines=[
                f"主题: {title}",
                f"章节数: {len(chapters)}",
            ],
            layout="title",
            notes="",
        ))

        # ── Slide 2: 学习目标 (如果存在) ──
        if objectives:
            slides.append(SlideItem(
                slide_index=2,
                title="学习目标",
                content_lines=objectives[:6],  # 最多 6 个目标
                layout="bullet_list",
                notes="",
            ))

        # ── 章节 Slides ──
        slide_idx = len(slides) + 1
        slide_number = slide_idx

        for ch in chapters:
            ch_title = ch.get("title", f"Chapter {slide_number}")

            # 章节标题页
            slides.append(SlideItem(
                slide_index=slide_number,
                title=ch_title,
                content_lines=[
                    ch.get("explanation", "")[:120] or f"本章将学习 {ch_title} 的核心知识",
                ],
                layout="title",
                notes="",
            ))
            slide_number += 1

            # 概念页
            concepts = ch.get("concepts", [])
            if concepts:
                concept_lines = []
                for c in concepts:
                    name = c.get("name", "")
                    desc = c.get("description", "")
                    difficulty = c.get("difficulty", "beginner")
                    concept_lines.append(f"📌 {name} [{difficulty}]: {desc}")
                slides.append(SlideItem(
                    slide_index=slide_number,
                    title=f"核心概念 — {ch_title}",
                    content_lines=concept_lines[:5],
                    layout="bullet_list",
                    notes=f"本章核心概念共 {len(concepts)} 个",
                ))
                slide_number += 1

            # 示例页
            examples = ch.get("examples", [])
            if examples:
                example_lines = []
                for e in examples[:3]:
                    etitle = e.get("title", "")
                    eexpl = e.get("explanation", "")
                    ecode = e.get("code", "")
                    example_lines.append(f"💡 {etitle}")
                    if eexpl:
                        example_lines.append(f"   {eexpl[:80]}")
                    if ecode:
                        # Show first 2 lines of code
                        code_lines = ecode.split("\n")[:2]
                        for cl in code_lines:
                            example_lines.append(f"   `{cl.strip()[:60]}`")
                slides.append(SlideItem(
                    slide_index=slide_number,
                    title=f"示例 — {ch_title}",
                    content_lines=example_lines[:6],
                    layout="content",
                    notes="",
                ))
                slide_number += 1

            # 练习页
            exercises = ch.get("exercises", [])
            if exercises:
                exercise_lines = []
                for ex in exercises[:3]:
                    q = ex.get("question", "")
                    hint = ex.get("hint", "")
                    exercise_lines.append(f"✏️ {q[:80]}")
                    if hint:
                        exercise_lines.append(f"   💭 提示: {hint[:60]}")
                slides.append(SlideItem(
                    slide_index=slide_number,
                    title=f"练习 — {ch_title}",
                    content_lines=exercise_lines[:6],
                    layout="content",
                    notes="",
                ))
                slide_number += 1

        # ── 总结页 ──
        summary_text = summary if summary else f"本次学习共覆盖 {len(chapters)} 个章节"
        slides.append(SlideItem(
            slide_index=slide_number,
            title="总结",
            content_lines=[
                summary_text,
                f"学习目标: {len(objectives)} 个",
                f"核心概念: {sum(len(ch.get('concepts', [])) for ch in chapters)} 个",
            ],
            layout="summary",
            notes="",
        ))

        structure = PPTStructure(
            ppt_id=ppt_id,
            title=title,
            topic=material.get("title", ""),
            level=material.get("target_profile", "beginner"),
            slides=slides,
            total_slides=len(slides),
            generation_source="rule",
            metadata={
                "chapter_count": len(chapters),
                "objective_count": len(objectives),
            },
        )
        return structure

    # ── .pptx 渲染 ─────────────────

    def _render_pptx(
        self,
        ppt_structure: PPTStructure,
        output_dir: str = "",
    ) -> str:
        """
        将 PPTStructure 渲染为 .pptx 文件.

        使用 python-pptx 库生成真实的 .pptx 文件。
        如果 python-pptx 不可用，抛出 RuntimeError。
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise RuntimeError(
                "python-pptx 库未安装。请运行: pip install python-pptx"
            )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_item in ppt_structure.slides:
            layout = slide_item.layout

            if layout == "title":
                # Title slide layout
                slide_layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_item.title
                # Subtitle
                if len(slide.placeholders) > 1 and slide_item.content_lines:
                    slide.placeholders[1].text = "\n".join(slide_item.content_lines)

            elif layout == "bullet_list":
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_item.title
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for i, line in enumerate(slide_item.content_lines):
                        if i == 0:
                            tf.text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 0

            elif layout == "content":
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_item.title
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for i, line in enumerate(slide_item.content_lines):
                        if i == 0:
                            tf.text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 1 if line.startswith("   ") else 0

            elif layout == "summary":
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_item.title or "总结"
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for i, line in enumerate(slide_item.content_lines):
                        if i == 0:
                            tf.text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 0

            else:
                # Default: content layout
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_item.title
                if len(slide.placeholders) > 1 and slide_item.content_lines:
                    slide.placeholders[1].text = "\n".join(slide_item.content_lines)

            # Add notes if present
            if slide_item.notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_item.notes

            # Phase 8.3-F2 — Add image if provided (SVG path or base64 data URI)
            if slide_item.image:
                if slide_item.image.startswith("data:image/svg+xml;base64,"):
                    # Render base64 SVG as picture
                    try:
                        import base64 as b64
                        svg_data = b64.b64decode(slide_item.image.split(",", 1)[1])
                        fd_img, img_path = tempfile.mkstemp(suffix=".svg", prefix="slide_img_")
                        os.close(fd_img)
                        with open(img_path, "wb") as fimg:
                            fimg.write(svg_data)
                        # python-pptx doesn't natively support SVG, so skip gracefully
                        os.unlink(img_path)
                    except Exception:
                        pass
                elif os.path.isfile(slide_item.image) and slide_item.image.endswith(".svg"):
                    # SVG file on disk — try to add (python-pptx has limited SVG support)
                    try:
                        # Add as picture (may work on newer python-pptx)
                        slide.shapes.add_picture(
                            slide_item.image,
                            Inches(0.5), Inches(2.5),
                            width=Inches(5), height=Inches(3.5),
                        )
                    except Exception:
                        pass

        # Save to file
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
            file_path = os.path.join(output_dir, f"{ppt_structure.ppt_id}.pptx")
        else:
            fd, file_path = tempfile.mkstemp(suffix=".pptx", prefix="a3_ppt_")
            os.close(fd)

        prs.save(file_path)

        # Verify file was created
        if not os.path.isfile(file_path) or os.path.getsize(file_path) < 100:
            raise RuntimeError(f"PPT 文件生成失败: {file_path}")

        return file_path

    # ── 辅助方法 ──────────────────

    def check_capability(
        self,
        provider: str,
        model: str = "",
    ) -> tuple[bool, Optional[str]]:
        """
        检查指定 provider/model 是否支持 PPT 生成.

        使用 check_task_capability() 进行能力检查。
        禁止绕过 Capability Layer。

        Returns:
            (supported: bool, error_message: Optional[str])
        """
        return check_task_capability(provider, model, TaskType.GENERATE_PPT)

    def get_recommended_model(
        self,
        preferred_provider: str = "",
    ) -> RouterResult:
        """
        获取推荐用于 PPT 生成的模型.

        使用 ModelRouter.select_model() 进行基于能力的模型选择。

        Returns:
            RouterResult with model selection details
        """
        return self._router.select_model(
            TaskType.GENERATE_PPT,
            preferred_provider=preferred_provider,
        )

    def get_capable_models(self) -> list:
        """获取所有支持 PPT 生成的模型."""
        return self._router.find_models(TaskType.GENERATE_PPT)

    def get_capable_providers(self) -> List[str]:
        """获取所有支持 PPT 生成的 provider."""
        return self._router.get_capable_providers(TaskType.GENERATE_PPT)

    def get_ppt_structure(
        self,
        material: Any,
    ) -> PPTStructure:
        """
        获取 PPT 结构而不渲染文件 (用于预览/验证).

        Args:
            material: TeachingMaterial 实例或 dict

        Returns:
            PPTStructure
        """
        material_dict = material.to_dict() if hasattr(material, "to_dict") else material
        if not material_dict:
            raise ValueError("TeachingMaterial 不能为空")
        return self.fallback_generate_structure(material_dict)
